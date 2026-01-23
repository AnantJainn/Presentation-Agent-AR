# agents/pptx_agent.py
from pptx import Presentation

def pptx_agent(state):
    prs = Presentation()

    for slide in state["presentation"]["slides"]:
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide["title"]

        tf = s.placeholders[1].text_frame
        tf.clear()
        for b in slide["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1

    prs.save("output.pptx")
    return state
